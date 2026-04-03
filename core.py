import os
import subprocess
from datetime import datetime
import zipfile

try:
    from spire.presentation import Presentation, FileFormat
except ImportError:
    Presentation = None

try:
    from pypdf import PdfWriter
except ImportError:
    PdfWriter = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None


def get_timestamp():
    return datetime.now().strftime("%Y%md_%H%M%S")

def merge_presentations(input_files, output_dir):
    """Merges multiple PPT or PPTX files into one."""
    if Presentation is None:
        raise RuntimeError("Spire.Presentation is not installed.")
    
    if not input_files:
        raise ValueError("No files provided for merging.")
        
    target_pres = Presentation()
    target_pres.LoadFromFile(input_files[0])
    
    for file_path in input_files[1:]:
        temp_pres = Presentation()
        temp_pres.LoadFromFile(file_path)
        for i in range(temp_pres.Slides.Count):
            target_pres.Slides.AppendBySlide(temp_pres.Slides[i])
        temp_pres.Dispose()
        
    output_filename = f"Merged_Presentations_{get_timestamp()}.pptx"
    output_path = os.path.join(output_dir, output_filename)
    target_pres.SaveToFile(output_path, FileFormat.Pptx2013)
    target_pres.Dispose()
    
    # Remove the Spire Evaluation Warning slide
    if PptxPresentation is not None:
        try:
            pptx_prs = PptxPresentation(output_path)
            slides_to_remove = []
            for slide in pptx_prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.lower()
                        if "evaluation" in text and "spire" in text:
                            slides_to_remove.append(slide._element)
                            break
            
            for sld in slides_to_remove:
                pptx_prs.slides._sldIdLst.remove(sld)
            
            pptx_prs.save(output_path)
        except Exception as e:
            print(f"Warning cleanup failed: {str(e)}")
    
    return output_path

def convert_presentations_to_pdf(input_files, output_dir):
    """Converts PPT/PPTX files to PDF. Zips them if multiple."""
    if not input_files:
        raise ValueError("No files provided for conversion.")
        
    generated_pdfs = []
    for file_path in input_files:
        result = subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            file_path,
            "--outdir", output_dir
        ], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        
        if result.returncode != 0:
            raise RuntimeError(f"Error converting {file_path}: {result.stderr}")
            
        # LibreOffice outputs the PDF with the same base name
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
        if os.path.exists(pdf_path):
            generated_pdfs.append(pdf_path)
        else:
            raise RuntimeError(f"Expected PDF not found at {pdf_path} after conversion.")

    if len(generated_pdfs) == 1:
        return generated_pdfs[0]
        
    # Zip if multiple
    zip_filename = f"Converted_PDFs_{get_timestamp()}.zip"
    zip_path = os.path.join(output_dir, zip_filename)
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        for pdf_file in generated_pdfs:
            zipf.write(pdf_file, arcname=os.path.basename(pdf_file))
            
    return zip_path

def merge_pdfs(input_files, output_dir):
    """Merges multiple PDF files into one."""
    if PdfWriter is None:
        raise RuntimeError("pypdf is not installed.")
        
    if not input_files:
        raise ValueError("No files provided for merging.")
        
    merger = PdfWriter()
    for pdf_file in input_files:
        merger.append(pdf_file)
        
    output_filename = f"Merged_PDFs_{get_timestamp()}.pdf"
    output_path = os.path.join(output_dir, output_filename)
    
    merger.write(output_path)
    merger.close()
    
    return output_path
